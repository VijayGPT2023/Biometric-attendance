"""Generate NPC Biometric Attendance System PPT presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(0x1a, 0x23, 0x7e)
NAVY_LIGHT = RGBColor(0x39, 0x49, 0xab)
ORANGE = RGBColor(0xe6, 0x51, 0x00)
GREEN = RGBColor(0x2e, 0x7d, 0x32)
RED = RGBColor(0xc6, 0x28, 0x28)
WHITE = RGBColor(0xff, 0xff, 0xff)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xf0, 0xf2, 0xf5)
BLACK = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2


def add_text_box(slide, left, top, width, height, text, size=18, color=BLACK,
                 bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_shape_box(slide, left, top, width, height, fill_color, text='',
                  font_size=11, font_color=WHITE, bold=False, radius=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Segoe UI'
    return shape


def add_arrow(slide, left, top):
    add_text_box(slide, left, top, 0.4, 0.4, '→', size=22, color=RGBColor(0x99, 0x99, 0x99),
                 align=PP_ALIGN.CENTER)


def add_icon_card(slide, left, top, icon, value, label, accent=NAVY):
    add_shape_box(slide, left, top, 2.0, 1.6, WHITE)
    add_text_box(slide, left, top + 0.1, 2.0, 0.5, icon, size=28, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.65, 2.0, 0.4, value, size=16, bold=True,
                 color=accent, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 1.05, 2.0, 0.4, label, size=9, color=GREY,
                 align=PP_ALIGN.CENTER)


def add_slide_number(slide, num, color=WHITE):
    add_text_box(slide, 12.5, 7.0, 0.7, 0.3, str(num), size=9,
                 color=color, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_text_box(s, 1, 1.2, 11, 0.5, 'NATIONAL PRODUCTIVITY COUNCIL', size=14,
             color=RGBColor(0xaa, 0xaa, 0xff), align=PP_ALIGN.CENTER)
add_text_box(s, 1, 1.8, 11, 1.2, 'Biometric Attendance\nManagement System', size=42,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.5, 11, 0.5, 'Ministry of Commerce & Industry, Government of India',
             size=14, color=RGBColor(0xcc, 0xcc, 0xff), align=PP_ALIGN.CENTER)

# Role badges
for i, (label, clr) in enumerate([
    ('DG / DDG', NAVY_LIGHT), ('Group Heads', ORANGE),
    ('Employees', GREEN), ('Admin', RED)
]):
    add_shape_box(s, 3.5 + i * 1.8, 4.5, 1.5, 0.45, clr, label, font_size=11, bold=True)

add_slide_number(s, 1)


# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape_box(s, 0.8, 0.7, 0.5, 0.06, NAVY)  # accent line
add_text_box(s, 0.8, 0.9, 6, 0.6, 'The Problem We Solve', size=32, bold=True, color=NAVY)

for i, (icon, val, desc) in enumerate([
    ('📋', 'Manual', 'Excel tracking\nis error-prone'),
    ('⏰', 'Delayed', 'Months to finalize\nattendance'),
    ('🚫', 'No Audit Trail', 'No record of\nwho approved what'),
    ('👤', 'Opaque', 'Employees can\'t\nsee own data'),
]):
    x = 1.0 + i * 2.8
    add_shape_box(s, x, 2.0, 2.3, 2.0, LIGHT_BG)
    add_text_box(s, x, 2.1, 2.3, 0.5, icon, size=30, align=PP_ALIGN.CENTER)
    add_text_box(s, x, 2.65, 2.3, 0.4, val, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, x, 3.1, 2.3, 0.7, desc, size=10, color=GREY, align=PP_ALIGN.CENTER)

# Before -> After
add_shape_box(s, 2.0, 4.8, 2.5, 0.6, RGBColor(0xff, 0xcd, 0xd2), '📋  Manual Excel',
              font_size=12, font_color=RED)
add_text_box(s, 4.7, 4.85, 0.5, 0.4, '→', size=24, color=GREY, align=PP_ALIGN.CENTER)
add_shape_box(s, 5.3, 4.8, 2.5, 0.6, RGBColor(0xe3, 0xf2, 0xfd), '⚡  This System',
              font_size=12, font_color=NAVY, bold=True)
add_text_box(s, 8.0, 4.85, 0.5, 0.4, '→', size=24, color=GREY, align=PP_ALIGN.CENTER)
add_shape_box(s, 8.6, 4.8, 2.8, 0.6, RGBColor(0xc8, 0xe6, 0xc9), '✅  Automated & Transparent',
              font_size=12, font_color=GREEN)

add_slide_number(s, 2, BLACK)


# ============================================================
# SLIDE 3: HOW IT WORKS (FLOW)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, RGBColor(0x66, 0x7e, 0xea))
add_text_box(s, 0.8, 0.8, 6, 0.6, 'How It Works', size=32, bold=True, color=WHITE)

steps = ['📁\nAdmin Uploads\nBiometric XLS', '⚠️\nSystem Detects\nAnomalies',
         '✍️\nEmployee\nJustifies', '🔍\nHead\nReviews',
         '✅\nAdmin\nFinalizes', '📅\nLeave\nDeducted']
for i, step in enumerate(steps):
    x = 0.6 + i * 2.1
    add_shape_box(s, x, 1.8, 1.7, 1.5, RGBColor(0x28, 0x35, 0x93), step,
                  font_size=10, font_color=WHITE)
    if i < 5:
        add_text_box(s, x + 1.75, 2.2, 0.35, 0.5, '→', size=20,
                     color=RGBColor(0x66, 0x7e, 0xea), align=PP_ALIGN.CENTER)

# Stats
for i, (val, label) in enumerate([('13+', 'Offices'), ('500+', 'Employees'),
                                    ('49', 'API Routes'), ('42', 'Tests Passing')]):
    x = 1.5 + i * 2.8
    add_shape_box(s, x, 4.2, 2.2, 1.5, WHITE)
    add_text_box(s, x, 4.35, 2.2, 0.7, val, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, x, 5.0, 2.2, 0.4, label, size=11, color=GREY, align=PP_ALIGN.CENTER)

add_slide_number(s, 3)


# ============================================================
# SLIDE 4: DETAILED WORKFLOW - TOP HALF
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape_box(s, 0.8, 0.4, 0.5, 0.06, NAVY)
add_text_box(s, 0.8, 0.55, 8, 0.5, 'Complete Workflow: Step by Step', size=28, bold=True, color=NAVY)

# STEP 1: Admin Upload
add_shape_box(s, 0.5, 1.4, 2.8, 2.4, RGBColor(0xfc, 0xe4, 0xec))
add_text_box(s, 0.5, 1.45, 2.8, 0.35, '❶  ADMIN UPLOADS', size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text_box(s, 0.7, 1.9, 2.4, 1.8,
    '• Selects Office (HQ / Regional)\n'
    '• Uploads biometric .xls file\n'
    '• Sets rules (late time, min hours)\n'
    '• System auto-detects anomalies\n'
    '• Creates record for each employee',
    size=10, color=BLACK)

# Arrow
add_text_box(s, 3.4, 2.2, 0.5, 0.5, '▶', size=24, color=RGBColor(0xcc, 0xcc, 0xcc), align=PP_ALIGN.CENTER)

# STEP 2: Employee Reviews & Justifies
add_shape_box(s, 4.0, 1.4, 3.6, 2.4, RGBColor(0xe8, 0xf5, 0xe9))
add_text_box(s, 4.0, 1.45, 3.6, 0.35, '❷  EMPLOYEE JUSTIFIES', size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text_box(s, 4.2, 1.9, 3.2, 1.8,
    '• Logs in → sees own anomalies\n'
    '• Each anomaly: Late / Early / Short Hrs\n'
    '• Writes reason for each anomaly\n'
    '   (tour, medical, device fault, etc.)\n'
    '• Clicks "Submit" → goes to Head\n'
    '• Can save draft & come back later',
    size=10, color=BLACK)

# Arrow
add_text_box(s, 7.7, 2.2, 0.5, 0.5, '▶', size=24, color=RGBColor(0xcc, 0xcc, 0xcc), align=PP_ALIGN.CENTER)

# STEP 3: Head Reviews
add_shape_box(s, 8.3, 1.4, 4.2, 2.4, RGBColor(0xff, 0xf3, 0xe0))
add_text_box(s, 8.3, 1.45, 4.2, 0.35, '❸  HEAD REVIEWS', size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text_box(s, 8.5, 1.9, 3.8, 0.6,
    '• Reads each justification\n'
    '• For EACH anomaly, decides:',
    size=10, color=BLACK)
# Three decision boxes
add_shape_box(s, 8.6, 2.65, 1.1, 0.55, GREEN, '✅ Accept\nExclude it', font_size=8, font_color=WHITE, bold=True)
add_shape_box(s, 9.8, 2.65, 1.1, 0.55, RED, '❌ Decline\nCount it', font_size=8, font_color=WHITE, bold=True)
add_shape_box(s, 11.0, 2.65, 1.3, 0.55, ORANGE, '❓ Query\nAsk more info', font_size=8, font_color=WHITE, bold=True)
add_text_box(s, 8.5, 3.3, 3.8, 0.4, '• Query → goes back to employee', size=9, color=GREY)

# Bottom: Step 4 Admin Finalize + Deduction
add_shape_box(s, 0.5, 4.2, 6.5, 2.6, RGBColor(0xe8, 0xea, 0xf6))
add_text_box(s, 0.5, 4.25, 6.5, 0.35, '❹  ADMIN FINALIZES', size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text_box(s, 0.7, 4.7, 6.1, 1.8,
    '• Reviews Head\'s decisions for all employees\n'
    '• Can override Head decision with recorded reason\n'
    '• For each anomaly:\n'
    '      "Exclude" → anomaly removed (no deduction)\n'
    '      "Count" → anomaly stays (leave deducted)\n'
    '• Clicks "Finalize Decisions" → month is locked',
    size=10, color=BLACK)

# Arrow to deduction
add_text_box(s, 7.2, 5.1, 0.5, 0.5, '▶', size=24, color=RGBColor(0xcc, 0xcc, 0xcc), align=PP_ALIGN.CENTER)

# Step 5: Leave Deduction
add_shape_box(s, 7.8, 4.2, 4.7, 2.6, RGBColor(0xff, 0xeb, 0xee))
add_text_box(s, 7.8, 4.25, 4.7, 0.35, '❺  LEAVE DEDUCTED', size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text_box(s, 8.0, 4.7, 4.3, 0.8,
    '• Only "Counted" anomalies remain\n'
    '• First 2 anomaly days = FREE\n'
    '• Each extra day = 0.5 EL deducted',
    size=10, color=BLACK)

# Example box
add_shape_box(s, 8.2, 5.6, 4.0, 1.0, WHITE)
add_text_box(s, 8.3, 5.65, 3.8, 0.8,
    'Example: 8 anomalies → 3 accepted\n'
    '= 5 counted − 2 free = 3 × 0.5 = 1.5 EL',
    size=10, bold=True, color=NAVY)

add_slide_number(s, 4, BLACK)


# ============================================================
# SLIDE 5: WORKFLOW VISUAL DIAGRAM
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_text_box(s, 0.8, 0.4, 8, 0.5, 'Workflow at a Glance', size=28, bold=True, color=WHITE)

# Large visual flow with swimlanes
# Admin lane
add_shape_box(s, 0.5, 1.2, 1.8, 0.45, RED, 'ADMIN', font_size=11, bold=True)
add_shape_box(s, 2.5, 1.2, 2.0, 0.45, RGBColor(0xff,0xcd,0xd2), 'Upload XLS', font_size=10, font_color=RED)
add_text_box(s, 4.6, 1.25, 0.4, 0.35, '→', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 5.0, 1.2, 2.0, 0.45, RGBColor(0xff,0xcd,0xd2), 'Auto-Analyze', font_size=10, font_color=RED)

# Employee lane
add_shape_box(s, 0.5, 2.2, 1.8, 0.45, GREEN, 'EMPLOYEE', font_size=11, bold=True)
add_text_box(s, 7.1, 1.25, 0.4, 0.35, '↓', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 2.5, 2.2, 2.0, 0.45, RGBColor(0xc8,0xe6,0xc9), 'View Anomalies', font_size=10, font_color=GREEN)
add_text_box(s, 4.6, 2.25, 0.4, 0.35, '→', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 5.0, 2.2, 2.0, 0.45, RGBColor(0xc8,0xe6,0xc9), 'Write Reasons', font_size=10, font_color=GREEN)
add_text_box(s, 7.1, 2.25, 0.4, 0.35, '→', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 7.5, 2.2, 1.8, 0.45, RGBColor(0xc8,0xe6,0xc9), 'Submit', font_size=10, font_color=GREEN, bold=True)

# Head lane
add_shape_box(s, 0.5, 3.2, 1.8, 0.45, ORANGE, 'HEAD', font_size=11, bold=True)
add_text_box(s, 9.4, 2.25, 0.4, 0.35, '↓', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 2.5, 3.2, 2.2, 0.45, RGBColor(0xff,0xe0,0xb2), 'Read Justification', font_size=10, font_color=ORANGE)
add_text_box(s, 4.8, 3.25, 0.4, 0.35, '→', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)

# Three outcomes
add_shape_box(s, 5.3, 3.0, 1.4, 0.4, GREEN, '✅ Accept', font_size=9, bold=True)
add_shape_box(s, 5.3, 3.5, 1.4, 0.4, RED, '❌ Decline', font_size=9, bold=True)
add_shape_box(s, 6.9, 3.2, 1.5, 0.45, ORANGE, '❓ Query', font_size=9, bold=True)

# Query loops back
add_text_box(s, 8.5, 3.25, 2.5, 0.35, '↩ back to employee', size=9, color=RGBColor(0xff,0xcc,0x80))

# Admin finalize lane
add_shape_box(s, 0.5, 4.2, 1.8, 0.45, RED, 'ADMIN', font_size=11, bold=True)
add_text_box(s, 6.8, 3.65, 0.4, 0.35, '↓', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 2.5, 4.2, 2.2, 0.45, RGBColor(0xff,0xcd,0xd2), 'Review Decisions', font_size=10, font_color=RED)
add_text_box(s, 4.8, 4.25, 0.4, 0.35, '→', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 5.3, 4.2, 1.8, 0.45, RGBColor(0xff,0xcd,0xd2), 'Finalize', font_size=10, font_color=RED, bold=True)
add_text_box(s, 7.2, 4.25, 0.4, 0.35, '→', size=16, color=RGBColor(0x66,0x7e,0xea), align=PP_ALIGN.CENTER)
add_shape_box(s, 7.6, 4.2, 2.0, 0.45, RGBColor(0xe8,0xea,0xf6), 'Deduct Leave', font_size=10, font_color=NAVY, bold=True)

# Deduction summary
add_shape_box(s, 1.0, 5.2, 11.3, 1.6, RGBColor(0x1a, 0x23, 0x7e))
add_text_box(s, 1.0, 5.25, 11.3, 0.35, 'Leave Deduction Summary', size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

labels = ['Total\nAnomalies', '−', 'Accepted\n(Excluded)', '−', 'Free\nAllowance (2)', '=', 'Deductible\nDays', '×', '0.5 EL\nper day', '=', 'Leave\nDeducted']
for i, lbl in enumerate(labels):
    x = 1.2 + i * 0.95
    if lbl in ('−', '=', '×'):
        add_text_box(s, x, 5.7, 0.7, 0.6, lbl, size=20, bold=True,
                     color=RGBColor(0xff, 0xcc, 0x80), align=PP_ALIGN.CENTER)
    else:
        clr = ORANGE if 'Deduct' in lbl else WHITE
        add_shape_box(s, x, 5.75, 0.85, 0.7,
                      RGBColor(0x28, 0x35, 0x93), lbl, font_size=8, font_color=clr, bold='Deduct' in lbl)

add_slide_number(s, 5)


# ============================================================
# SLIDE 6: LOGIN (was slide 4)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, NAVY)
add_text_box(s, 0.8, 0.8, 5, 0.6, 'Login Guide', size=32, bold=True, color=NAVY)

# Login mockup
add_shape_box(s, 1.0, 1.8, 5, 4.5, WHITE)
add_text_box(s, 1.0, 1.9, 5, 0.4, '🌐  npc-biometric-attendance.up.railway.app', size=10,
             color=GREY, align=PP_ALIGN.CENTER)
add_text_box(s, 1.0, 2.5, 5, 0.5, 'NPC Biometric Attendance', size=18, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
add_text_box(s, 1.0, 3.0, 5, 0.3, 'राष्ट्रीय उत्पादकता परिषद', size=11,
             color=GREY, align=PP_ALIGN.CENTER)
add_shape_box(s, 2.0, 3.6, 3, 0.45, LIGHT_BG, 'Username: firstname.lastname',
              font_size=10, font_color=GREY)
add_shape_box(s, 2.0, 4.2, 3, 0.45, LIGHT_BG, 'Password: ********',
              font_size=10, font_color=GREY)
add_shape_box(s, 2.0, 4.9, 3, 0.5, NAVY, 'Sign In', font_size=14, bold=True)
add_text_box(s, 2.0, 5.55, 3, 0.3, 'Forgot Password?', size=10, color=NAVY, align=PP_ALIGN.CENTER)

# Credentials
add_text_box(s, 7.0, 1.8, 5, 0.4, '👤  Username Format', size=16, bold=True, color=NAVY)
creds = [
    ('Employee', 'dk.rahul', ''),
    ('Head', 'gh.hrmgroup', ''),
    ('DG', 'dg', ''),
    ('DDG-I', 'ddg1', ''),
    ('DDG-II', 'ddg2', ''),
]
for i, (role, uname, _) in enumerate(creds):
    y = 2.4 + i * 0.4
    add_text_box(s, 7.2, y, 2, 0.35, role, size=11, color=GREY)
    add_shape_box(s, 9.3, y, 2.2, 0.32, RGBColor(0xe8, 0xea, 0xf6), uname,
                  font_size=11, font_color=NAVY, bold=True)

add_text_box(s, 7.0, 4.6, 5, 0.4, '🔐  First Login', size=16, bold=True, color=NAVY)
add_text_box(s, 7.2, 5.1, 5, 0.3, '• Default password: npc123', size=12, color=BLACK)
add_text_box(s, 7.2, 5.45, 5, 0.3, '• Must change immediately', size=12, color=RED)
add_text_box(s, 7.2, 5.8, 5, 0.3, '• Min 8 chars + letter + number', size=12, color=BLACK)

add_slide_number(s, 4, BLACK)


# ============================================================
# SLIDE 5: SECTION A - DG/DDG
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_text_box(s, 1, 1.5, 11, 0.4, 'SECTION A', size=14,
             color=RGBColor(0xaa, 0xaa, 0xff), align=PP_ALIGN.CENTER)
add_text_box(s, 1, 2.2, 11, 1.0, '🏢  DG & DDG', size=44,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.8, 11, 0.5, 'Organization-wide attendance oversight',
             size=18, color=RGBColor(0xcc, 0xcc, 0xff), align=PP_ALIGN.CENTER)
add_slide_number(s, 5)


# ============================================================
# SLIDE 6: DG VIEW
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, NAVY)
add_text_box(s, 0.8, 0.8, 6, 0.6, 'DG / DDG: Your View', size=28, bold=True, color=NAVY)

# Department tabs visual
tabs = ['AB Group', 'Admin', 'ECA Group', 'Finance', 'HRM', 'IE Group']
for i, tab in enumerate(tabs):
    clr = WHITE if i != 2 else NAVY
    bg = RGBColor(0xff, 0xf3, 0xe0) if i != 2 else WHITE
    fc = ORANGE if i != 2 else ORANGE
    add_shape_box(s, 1.0 + i * 1.8, 1.8, 1.5, 0.45, bg, tab, font_size=9, font_color=fc, bold=True)

add_text_box(s, 1.0, 2.5, 10, 0.3, 'Click any department tab → see employee summary → click row → daily detail',
             size=11, color=GREY)

# What you see / do
add_shape_box(s, 1.0, 3.2, 5.2, 3.2, LIGHT_BG)
add_text_box(s, 1.2, 3.3, 5, 0.35, '👁️  What You See', size=14, bold=True, color=NAVY)
items = ['All 12 departments as clickable tabs',
         'Employee summary per department',
         'Anomaly counts & leave deductions',
         'Justification status at a glance']
for i, item in enumerate(items):
    add_text_box(s, 1.4, 3.8 + i * 0.45, 4.5, 0.35, '✓  ' + item, size=11, color=BLACK)

add_shape_box(s, 6.8, 3.2, 5.5, 3.2, RGBColor(0xe8, 0xf5, 0xe9))
add_text_box(s, 7.0, 3.3, 5, 0.35, '👆  What You Do', size=14, bold=True, color=GREEN)
items2 = ['Click department → see employees',
          'Click employee → see daily anomalies',
          'Monitor chronic late-comers',
          'View justification progress',
          'No action required (oversight only)']
for i, item in enumerate(items2):
    add_text_box(s, 7.2, 3.8 + i * 0.45, 5, 0.35, '✓  ' + item, size=11, color=BLACK)

add_slide_number(s, 6, BLACK)


# ============================================================
# SLIDE 7: SECTION B - GROUP HEADS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0xbf, 0x36, 0x0c), ORANGE)
add_text_box(s, 1, 1.5, 11, 0.4, 'SECTION B', size=14,
             color=RGBColor(0xff, 0xcc, 0x80), align=PP_ALIGN.CENTER)
add_text_box(s, 1, 2.2, 11, 1.0, '👥  Group Heads', size=44,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.8, 11, 0.5, 'Review & decide on employee justifications',
             size=18, color=RGBColor(0xff, 0xe0, 0xb2), align=PP_ALIGN.CENTER)
add_slide_number(s, 7)


# ============================================================
# SLIDE 8: HEAD WORKFLOW
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, ORANGE)
add_text_box(s, 0.8, 0.8, 6, 0.5, 'Group Head: Workflow', size=28, bold=True, color=ORANGE)

# Flow
flow_items = [
    ('📊\nDashboard', RGBColor(0xff, 0xe0, 0xb2), ORANGE),
    ('🔍\nReview', RGBColor(0xff, 0xe0, 0xb2), ORANGE),
    ('✅\nAccept', RGBColor(0xc8, 0xe6, 0xc9), GREEN),
    ('❌\nDecline', RGBColor(0xff, 0xcd, 0xd2), RED),
    ('❓\nQuery', RGBColor(0xff, 0xf3, 0xe0), ORANGE),
]
for i, (text, bg, fc) in enumerate(flow_items):
    x = 0.8 + i * 2.5
    add_shape_box(s, x, 1.7, 2.0, 1.2, bg, text, font_size=11, font_color=fc, bold=True)
    if i < 2:
        add_text_box(s, x + 2.05, 2.0, 0.4, 0.4, '→', size=20, color=GREY, align=PP_ALIGN.CENTER)

# Status pills explanation
add_shape_box(s, 0.8, 3.4, 6.0, 2.8, RGBColor(0xff, 0xf8, 0xe1))
add_text_box(s, 1.0, 3.5, 5, 0.35, 'Status Indicators', size=14, bold=True, color=ORANGE)
pills = [
    ('12P', RGBColor(0x9e, 0x9e, 0x9e), 'Pending - employee hasn\'t submitted'),
    ('5S', RGBColor(0x19, 0x76, 0xd2), 'Submitted - ready for your review'),
    ('2Q', ORANGE, 'Queried - sent back to employee'),
    ('8A', GREEN, 'Accepted - anomaly excluded'),
    ('3D', RED, 'Declined - counts for deduction'),
]
for i, (pill, clr, desc) in enumerate(pills):
    y = 4.0 + i * 0.4
    add_shape_box(s, 1.2, y, 0.6, 0.3, clr, pill, font_size=9, font_color=WHITE, bold=True)
    add_text_box(s, 2.0, y, 4.5, 0.3, desc, size=11, color=BLACK)

# Tips
add_shape_box(s, 7.3, 3.4, 5.0, 2.8, RGBColor(0xe8, 0xf5, 0xe9))
add_text_box(s, 7.5, 3.5, 4.5, 0.35, '💡  Tips', size=14, bold=True, color=GREEN)
tips = ['Use "Accept All" for bulk approvals',
        'Add remarks when declining',
        'Query sends notification to employee',
        'Click employee name for full detail',
        'Bulk actions save time']
for i, tip in enumerate(tips):
    add_text_box(s, 7.7, 4.0 + i * 0.4, 4.5, 0.35, '✓  ' + tip, size=11, color=BLACK)

add_slide_number(s, 8, BLACK)


# ============================================================
# SLIDE 9: SECTION C - EMPLOYEES
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x1b, 0x5e, 0x20), GREEN)
add_text_box(s, 1, 1.5, 11, 0.4, 'SECTION C', size=14,
             color=RGBColor(0xa5, 0xd6, 0xa7), align=PP_ALIGN.CENTER)
add_text_box(s, 1, 2.2, 11, 1.0, '👤  Employees', size=44,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.8, 11, 0.5, 'View your attendance & justify anomalies',
             size=18, color=RGBColor(0xc8, 0xe6, 0xc9), align=PP_ALIGN.CENTER)
add_slide_number(s, 9)


# ============================================================
# SLIDE 10: EMPLOYEE - 3 STEPS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, GREEN)
add_text_box(s, 0.8, 0.8, 6, 0.5, 'Employee: 3 Simple Steps', size=28, bold=True, color=GREEN)

# 3 big numbered steps
for i, (num, title, desc) in enumerate([
    ('1', 'View Anomalies', 'See your attendance summary\nand day-wise anomaly list'),
    ('2', 'Write Justification', 'Type your reason for each\nanomaly (tour, medical, etc.)'),
    ('3', 'Submit', 'Click Submit → sent to Head\nSuccess message with timestamp'),
]):
    x = 1.0 + i * 4.0
    # Number circle
    add_shape_box(s, x + 0.8, 1.7, 0.8, 0.8, GREEN, num, font_size=28, bold=True)
    add_text_box(s, x, 2.7, 3.2, 0.4, title, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text_box(s, x, 3.2, 3.2, 0.8, desc, size=12, color=GREY, align=PP_ALIGN.CENTER)

# Summary cards mockup
cards_data = [('58', 'Working\nDays', LIGHT_BG, NAVY), ('52', 'Present', RGBColor(0xe8,0xf5,0xe9), GREEN),
              ('6', 'Anomalies', RGBColor(0xff,0xeb,0xee), RED), ('2.0', 'Leave\nDed.', RGBColor(0xff,0xf3,0xe0), ORANGE)]
for i, (val, lbl, bg, fc) in enumerate(cards_data):
    x = 1.5 + i * 2.8
    add_shape_box(s, x, 4.5, 2.2, 1.4, bg)
    add_text_box(s, x, 4.6, 2.2, 0.6, val, size=28, bold=True, color=fc, align=PP_ALIGN.CENTER)
    add_text_box(s, x, 5.2, 2.2, 0.5, lbl, size=10, color=GREY, align=PP_ALIGN.CENTER)

# Offline note
add_shape_box(s, 1.0, 6.2, 11.3, 0.55, RGBColor(0xff, 0xf3, 0xe0),
              '🔌  Offline? Data auto-saved locally. Submits automatically when connection restores.',
              font_size=11, font_color=ORANGE)

add_slide_number(s, 10, BLACK)


# ============================================================
# SLIDE 11: LEAVE DEDUCTION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, ORANGE)
add_text_box(s, 0.8, 0.8, 4, 0.5, 'Leave Deduction', size=32, bold=True, color=WHITE)

# Big numbers
add_text_box(s, 1.5, 2.0, 3.5, 1.2, '2', size=80, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text_box(s, 1.5, 3.5, 3.5, 0.4, 'Free anomaly days\nper month', size=13,
             color=RGBColor(0xcc, 0xcc, 0xff), align=PP_ALIGN.CENTER)

add_text_box(s, 1.5, 4.5, 3.5, 0.8, '0.5', size=56, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text_box(s, 1.5, 5.5, 3.5, 0.4, 'EL deducted per\nextra anomaly day', size=13,
             color=RGBColor(0xcc, 0xcc, 0xff), align=PP_ALIGN.CENTER)

# Table
table_data = [('1-2 days', '0 EL', '🟢🟢'),
              ('3 days', '0.5 EL', '🟢🟢🔴'),
              ('5 days', '1.5 EL', '🟢🟢🔴🔴🔴'),
              ('8 days', '3.0 EL', '🟢🟢🔴🔴🔴🔴🔴🔴'),
              ('12 days', '5.0 EL', '🟢🟢🔴 ×10')]

add_shape_box(s, 6.0, 1.8, 6.2, 4.5, RGBColor(0x1a, 0x23, 0x7e))
add_text_box(s, 6.3, 1.9, 1.8, 0.4, 'Anomaly Days', size=11, bold=True,
             color=RGBColor(0xaa, 0xaa, 0xff))
add_text_box(s, 8.3, 1.9, 1.5, 0.4, 'Deduction', size=11, bold=True,
             color=RGBColor(0xaa, 0xaa, 0xff))
add_text_box(s, 10.0, 1.9, 2.0, 0.4, 'Visual', size=11, bold=True,
             color=RGBColor(0xaa, 0xaa, 0xff))

for i, (days, ded, viz) in enumerate(table_data):
    y = 2.5 + i * 0.7
    color = WHITE
    if '3.0' in ded or '5.0' in ded:
        color = ORANGE
    add_text_box(s, 6.3, y, 1.8, 0.4, days, size=13, color=WHITE)
    add_text_box(s, 8.3, y, 1.5, 0.4, ded, size=13, bold=True, color=color)
    add_text_box(s, 10.0, y, 2.0, 0.4, viz, size=12, color=WHITE)

add_text_box(s, 6.3, 6.1, 5.5, 0.4, '🟢 = Free    🔴 = 0.5 EL each    ✅ Accepted = excluded',
             size=10, color=RGBColor(0x99, 0x99, 0xcc))

add_slide_number(s, 11)


# ============================================================
# SLIDE 12: SECTION D - ADMIN
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x7f, 0x00, 0x00), RED)
add_text_box(s, 1, 1.5, 11, 0.4, 'SECTION D', size=14,
             color=RGBColor(0xff, 0x8a, 0x80), align=PP_ALIGN.CENTER)
add_text_box(s, 1, 2.2, 11, 1.0, '⚙️  Admin', size=44,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.8, 11, 0.5, 'Upload, manage, finalize, export',
             size=18, color=RGBColor(0xff, 0xcd, 0xd2), align=PP_ALIGN.CENTER)
add_slide_number(s, 12)


# ============================================================
# SLIDE 13: ADMIN PROCESS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, RED)
add_text_box(s, 0.8, 0.8, 6, 0.5, 'Admin: Monthly Process', size=28, bold=True, color=RED)

# Process flow
admin_steps = [
    ('📁\nUpload', RGBColor(0xff,0xcd,0xd2), RED),
    ('⚡\nAnalyze', RGBColor(0xe3,0xf2,0xfd), NAVY),
    ('⏳\nWait', RGBColor(0xff,0xf3,0xe0), ORANGE),
    ('✅\nFinalize', RGBColor(0xc8,0xe6,0xc9), GREEN),
    ('📄\nExport', RGBColor(0xe8,0xea,0xf6), NAVY),
]
for i, (text, bg, fc) in enumerate(admin_steps):
    x = 0.6 + i * 2.5
    add_shape_box(s, x, 1.6, 2.0, 1.1, bg, text, font_size=11, font_color=fc, bold=True)
    if i < 4:
        add_text_box(s, x + 2.05, 1.8, 0.4, 0.5, '→', size=20, color=GREY, align=PP_ALIGN.CENTER)

# Two columns
add_shape_box(s, 0.8, 3.2, 5.8, 3.5, RGBColor(0xfc, 0xe4, 0xec))
add_text_box(s, 1.0, 3.3, 5, 0.35, '🔧  Admin Tools', size=14, bold=True, color=RED)
tools = ['User management (add/edit/reset/delete)',
         'Office & department management',
         'Holiday calendar management',
         'Anomaly rules configuration',
         'Password reset for employees']
for i, tool in enumerate(tools):
    add_text_box(s, 1.2, 3.8 + i * 0.45, 5, 0.35, '✓  ' + tool, size=11, color=BLACK)

add_shape_box(s, 7.0, 3.2, 5.5, 3.5, RGBColor(0xe8, 0xea, 0xf6))
add_text_box(s, 7.2, 3.3, 5, 0.35, '📊  Reports & Data', size=14, bold=True, color=NAVY)
reports = ['Excel export (dept-wise sheets)',
           'PDF print-friendly report',
           'Audit log (who did what, when)',
           'eHRMS leave reconciliation',
           'Chart.js visual analytics']
for i, rep in enumerate(reports):
    add_text_box(s, 7.4, 3.8 + i * 0.45, 5, 0.35, '✓  ' + rep, size=11, color=BLACK)

add_slide_number(s, 13, BLACK)


# ============================================================
# SLIDE 14: SECURITY
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, ORANGE)
add_text_box(s, 0.8, 0.8, 6, 0.5, 'Security & Compliance', size=32, bold=True, color=WHITE)

features = [
    ('🔐', 'CSRF', 'All forms protected'),
    ('🔑', 'Password', 'Strength + forced change'),
    ('⏰', 'Timeout', '30 min auto-logout'),
    ('📝', 'Audit', 'Every action logged'),
    ('🌐', 'GIGW 3.0', 'Compliant UI'),
    ('📱', 'Mobile', 'Responsive design'),
    ('🇮🇳', 'Hindi', 'Bilingual ready'),
    ('🔌', 'Offline', 'Local data save'),
    ('📊', 'Charts', 'Visual analytics'),
    ('⏱️', 'Rate Limit', 'Brute-force protected'),
]
for i, (icon, val, desc) in enumerate(features):
    row = i // 5
    col = i % 5
    x = 0.8 + col * 2.4
    y = 1.8 + row * 2.4
    add_shape_box(s, x, y, 2.0, 1.8, WHITE)
    add_text_box(s, x, y + 0.15, 2.0, 0.5, icon, size=26, align=PP_ALIGN.CENTER)
    add_text_box(s, x, y + 0.7, 2.0, 0.4, val, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, x, y + 1.15, 2.0, 0.4, desc, size=9, color=GREY, align=PP_ALIGN.CENTER)

add_slide_number(s, 14)


# ============================================================
# SLIDE 15: TECH STACK
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_shape_box(s, 0.8, 0.6, 0.5, 0.06, NAVY)
add_text_box(s, 0.8, 0.8, 6, 0.5, 'Technology Stack', size=28, bold=True, color=NAVY)

stack = [('Backend', 'Python Flask (12 Blueprints)'),
         ('Database', 'PostgreSQL (18 tables)'),
         ('ORM', 'SQLAlchemy + Flask-Migrate'),
         ('Auth', 'Flask-Login + CSRF'),
         ('Frontend', 'Jinja2 + Chart.js'),
         ('Hosting', 'Railway (Cloud PaaS)'),
         ('Testing', 'Playwright (42 E2E tests)')]

add_shape_box(s, 0.8, 1.6, 7, 5.0, WHITE)
for i, (layer, tech) in enumerate(stack):
    y = 1.8 + i * 0.65
    add_text_box(s, 1.2, y, 2.5, 0.4, layer, size=12, bold=True, color=NAVY)
    add_text_box(s, 3.8, y, 4, 0.4, tech, size=12, color=BLACK)
    if i < 6:
        # Divider line
        add_shape_box(s, 1.2, y + 0.5, 6.2, 0.01, RGBColor(0xee, 0xee, 0xee))

# Big stats
for i, (val, lbl, clr) in enumerate([('49', 'API Routes', NAVY),
                                       ('18', 'DB Tables', GREEN),
                                       ('42', 'Tests', ORANGE)]):
    x = 8.5
    y = 1.8 + i * 1.8
    add_shape_box(s, x, y, 3.8, 1.4, RGBColor(0xe8, 0xea, 0xf6) if i == 0
                  else (RGBColor(0xe8, 0xf5, 0xe9) if i == 1 else RGBColor(0xff, 0xf3, 0xe0)))
    add_text_box(s, x, y + 0.1, 3.8, 0.7, val, size=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text_box(s, x, y + 0.85, 3.8, 0.3, lbl, size=11, color=GREY, align=PP_ALIGN.CENTER)

add_slide_number(s, 15, BLACK)


# ============================================================
# SLIDE 16: THANK YOU
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(s, RGBColor(0x0d, 0x1b, 0x4a), NAVY_LIGHT)
add_text_box(s, 1, 1.5, 11, 1.0, 'Thank You', size=48,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.0, 11, 0.5, 'NPC Biometric Attendance Management System',
             size=18, color=RGBColor(0xcc, 0xcc, 0xff), align=PP_ALIGN.CENTER)

add_shape_box(s, 3.5, 4.0, 6, 1.5, RGBColor(0x28, 0x35, 0x93))
add_text_box(s, 3.5, 4.2, 6, 0.4, '📧  Support: vijay.kumar@npcindia.gov.in',
             size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, 3.5, 4.7, 6, 0.4, '🌐  npc-biometric-attendance.up.railway.app',
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s, 1, 6.0, 11, 0.5,
             'National Productivity Council  •  Ministry of Commerce & Industry  •  Government of India',
             size=11, color=RGBColor(0x99, 0x99, 0xcc), align=PP_ALIGN.CENTER)
add_slide_number(s, 16)


# ============================================================
# SAVE
# ============================================================
output_path = 'docs/NPC_Biometric_v2.pptx'
prs.save(output_path)
print(f'PPT saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
